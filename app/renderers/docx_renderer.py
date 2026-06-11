from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE as PPTX_XL_CHART_TYPE
from pptx.util import Inches as PPTXInches
import zipfile
import os

from docx import Document as WordDocument

from app.models.document import Paragraph
from app.models.table import Table
from app.models.chart import Chart
from app.services.openxml_patcher import OXmlPatcher


class DocxRenderer:
    """Renders document models to DOCX files."""

    @staticmethod
    def render(document, output_path):

        doc = WordDocument()
        charts_to_inject = []

        for element in document.elements:

            # Paragraphs / Headings
            if isinstance(element, Paragraph):

                DocxRenderer._render_paragraph(
                    doc,
                    element
                )

            # Tables
            elif isinstance(element, Table):

                DocxRenderer._render_table(
                    doc,
                    element
                )

            # Charts
            elif isinstance(element, Chart):

                charts_to_inject.append(
                    element
                )

                doc.add_paragraph(
                    f"CHART_PLACEHOLDER_{element.chart_id}"
                )

        doc.save(output_path)

        if charts_to_inject:

            DocxRenderer._inject_native_charts(
                output_path,
                charts_to_inject
            )

    @staticmethod
    def _render_paragraph(doc, paragraph):

        heading_levels = {
            "h1": 1,
            "h2": 2,
            "h3": 3,
            "h4": 4,
            "h5": 5,
            "h6": 6,
        }

        if paragraph.tag in heading_levels:

            doc.add_heading(
                paragraph.text,
                level=heading_levels[
                    paragraph.tag
                ]
            )

        else:

            doc.add_paragraph(
                paragraph.text
            )

    @staticmethod
    def _inject_native_charts(output_path, charts):

        if not charts:
            return

        chart_type_map = {
            "bar": PPTX_XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": PPTX_XL_CHART_TYPE.LINE,
            "pie": PPTX_XL_CHART_TYPE.PIE,
            "area": PPTX_XL_CHART_TYPE.AREA,
        }

        chart_files = []

        for idx, chart_model in enumerate(
            charts,
            start=1
        ):

            try:

                temp_pptx = (
                    f"temp_chart_gen_{idx}.pptx"
                )

                prs = Presentation()

                slide = prs.slides.add_slide(
                    prs.slide_layouts[6]
                )

                chart_data = CategoryChartData()

                chart_data.categories = (
                    chart_model.categories
                )

                for series in chart_model.series:

                    chart_data.add_series(
                        series.name,
                        tuple(series.values)
                    )

                x, y, cx, cy = (
                    PPTXInches(1),
                    PPTXInches(1),
                    PPTXInches(6),
                    PPTXInches(4),
                )

                chart_type = chart_type_map.get(
                    chart_model.chart_type,
                    PPTX_XL_CHART_TYPE.COLUMN_CLUSTERED
                )

                slide.shapes.add_chart(
                    chart_type,
                    x,
                    y,
                    cx,
                    cy,
                    chart_data
                )

                prs.save(temp_pptx)

                with zipfile.ZipFile(
                    temp_pptx,
                    "r"
                ) as z:

                    excel_sheet_path = None

                    for fname in z.namelist():

                        if fname.startswith(
                            "ppt/embeddings/Microsoft_Excel_"
                        ):
                            excel_sheet_path = fname
                            break

                    if not excel_sheet_path:

                        raise FileNotFoundError(
                            "No embedded workbook found"
                        )

                    xml = z.read(
                        "ppt/charts/chart1.xml"
                    )

                    rels = z.read(
                        "ppt/charts/_rels/chart1.xml.rels"
                    )

                    xlsx = z.read(
                        excel_sheet_path
                    )

                chart_files.append(
                    {
                        "idx": idx,
                        "chart_id": chart_model.chart_id,
                        "xml": xml,
                        "rels": rels,
                        "xlsx": xlsx,
                    }
                )

            except Exception as e:

                print(
                    f"Warning: Failed to generate chart {idx}: {e}"
                )

            finally:

                if os.path.exists(
                    temp_pptx
                ):

                    # Uncomment later
                    # os.remove(temp_pptx)
                    pass

        if chart_files:

            OXmlPatcher.inject_charts(
                output_path,
                chart_files
            )

    @staticmethod
    def _render_table(doc, table_model):

        table = doc.add_table(
            rows=len(table_model.rows) + 1,
            cols=len(table_model.headers)
        )

        table.style = "Table Grid"

        for col, header in enumerate(
            table_model.headers
        ):

            table.cell(
                0,
                col
            ).text = str(header)

        for row_idx, row in enumerate(
            table_model.rows
        ):

            for col_idx, value in enumerate(
                row
            ):

                table.cell(
                    row_idx + 1,
                    col_idx
                ).text = str(value)