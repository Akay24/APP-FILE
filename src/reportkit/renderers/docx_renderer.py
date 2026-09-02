"""DOCX renderer.

Converts a :class:`Document` model into a Microsoft Word ``.docx`` file,
including native Office charts injected via OpenXML patching.
"""

from __future__ import annotations

import logging
import os
import tempfile
import zipfile
from typing import Any

from docx import Document as WordDocument
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE as PPTX_XL_CHART_TYPE
from pptx.util import Inches as PPTXInches

from reportkit.exceptions import RenderingError
from reportkit.models.chart import Chart
from reportkit.models.document import Document, Paragraph
from reportkit.models.table import Table
from reportkit.services.openxml_patcher import OXmlPatcher

__all__ = ["DocxRenderer"]

logger = logging.getLogger(__name__)

_HEADING_LEVELS: dict[str, int] = {
    "h1": 1,
    "h2": 2,
    "h3": 3,
    "h4": 4,
    "h5": 5,
    "h6": 6,
}

_CHART_TYPE_MAP: dict[str, PPTX_XL_CHART_TYPE] = {
    "bar": PPTX_XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": PPTX_XL_CHART_TYPE.LINE,
    "pie": PPTX_XL_CHART_TYPE.PIE,
    "area": PPTX_XL_CHART_TYPE.AREA,
}


class DocxRenderer:
    """Renders a :class:`Document` to a ``.docx`` file.

    Charts are injected as native Office charts (not images) by:

    1. Generating a temporary PPTX with the chart via python-pptx.
    2. Extracting the chart XML and embedded Excel workbook.
    3. Patching them into the DOCX via :class:`OXmlPatcher`.
    """

    @staticmethod
    def render(document: Document, output_path: str) -> None:
        """Render a document model to a DOCX file.

        Args:
            document: Parsed document model.
            output_path: Filesystem path for the output ``.docx``.

        Raises:
            RenderingError: If chart injection fails.
        """
        doc = WordDocument()

        # Configure page layout (A4 with 0.75-inch margins) to match PDF
        section = doc.sections[0]
        section.page_width = Inches(8.27)
        section.page_height = Inches(11.69)
        section.left_margin = Inches(0.75)
        section.right_margin = Inches(0.75)
        section.top_margin = Inches(0.75)
        section.bottom_margin = Inches(0.75)

        # Set default font style (Arial, 10.5pt, charcoal gray) to match PDF
        normal_style = doc.styles["Normal"]
        font = normal_style.font
        font.name = "Arial"
        font.size = Pt(10.5)
        font.color.rgb = RGBColor(51, 51, 51)

        # Style Headings to match PDF's Helvetica styling
        for level in _HEADING_LEVELS.values():
            style_name = f"Heading {level}"
            if style_name in doc.styles:
                heading_style = doc.styles[style_name]
                h_font = heading_style.font
                h_font.name = "Arial"
                h_font.bold = True
                if level == 1:
                    h_font.size = Pt(22)
                    h_font.color.rgb = RGBColor(38, 68, 120)  # Hex #264478
                elif level == 2:
                    h_font.size = Pt(18)
                    h_font.color.rgb = RGBColor(68, 114, 196)  # Hex #4472C4
                elif level == 3:
                    h_font.size = Pt(14)
                    h_font.color.rgb = RGBColor(68, 114, 196)
                else:
                    h_font.size = Pt(12)
                    h_font.color.rgb = RGBColor(51, 51, 51)

        charts_to_inject: list[Chart] = []

        for element in document.elements:
            if isinstance(element, Paragraph):
                DocxRenderer._render_paragraph(doc, element)
            elif isinstance(element, Table):
                DocxRenderer._render_table(doc, element)
            elif isinstance(element, Chart):
                charts_to_inject.append(element)
                doc.add_paragraph(f"CHART_PLACEHOLDER_{element.chart_id}")

        doc.save(output_path)

        if charts_to_inject:
            DocxRenderer._inject_native_charts(output_path, charts_to_inject)

    @staticmethod
    def _render_paragraph(doc: Any, paragraph: Paragraph) -> None:
        """Add a heading or paragraph to the Word document."""
        level = _HEADING_LEVELS.get(paragraph.tag)
        if level is not None:
            doc.add_heading(paragraph.text, level=level)
        else:
            doc.add_paragraph(paragraph.text)

    @staticmethod
    def _render_table(doc: Any, table_model: Table) -> None:
        """Add a styled table to the Word document matching the PDF theme."""
        if not table_model.headers and not table_model.rows:
            return

        row_count = len(table_model.rows)
        if table_model.headers:
            row_count += 1
        col_count = (
            len(table_model.headers)
            if table_model.headers
            else (len(table_model.rows[0]) if table_model.rows else 0)
        )
        if col_count == 0:
            return

        table = doc.add_table(rows=row_count, cols=col_count)
        table.style = "Table Grid"

        current_row = 0

        # Render headers (White text on Medium Blue background)
        if table_model.headers:
            for col_idx, header in enumerate(table_model.headers):
                cell = table.cell(current_row, col_idx)
                DocxRenderer._set_cell_style(cell, header, is_header=True)
            current_row += 1

        # Render rows (Alternating light gray background)
        for row_idx, row in enumerate(table_model.rows):
            is_alternate = row_idx % 2 == 1
            for col_idx, value in enumerate(row):
                if col_idx < col_count:
                    cell = table.cell(current_row, col_idx)
                    DocxRenderer._set_cell_style(
                        cell, value, is_header=False, is_alternate=is_alternate
                    )
            current_row += 1

    @staticmethod
    def _set_cell_style(
        cell: Any,
        text: str,
        is_header: bool = False,
        is_alternate: bool = False,
    ) -> None:
        """Apply shading and font styling to a cell to match the PDF design system."""
        # 1. Shading
        if is_header:
            DocxRenderer._set_cell_shading(cell, "4472C4")  # Medium Blue (#4472C4)
        elif is_alternate:
            DocxRenderer._set_cell_shading(cell, "F2F2F2")  # Light Gray (#F2F2F2)

        # 2. Text styling & margins
        if not cell.paragraphs:
            p = cell.add_paragraph()
        else:
            p = cell.paragraphs[0]
            p.text = ""

        p.paragraph_format.space_before = Pt(4)
        p.paragraph_format.space_after = Pt(4)
        p.paragraph_format.line_spacing = 1.15

        run = p.add_run(text)
        run.font.name = "Arial"

        if is_header:
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)  # White
        else:
            run.font.size = Pt(9.5)
            run.font.bold = False
            run.font.color.rgb = RGBColor(51, 51, 51)  # Charcoal

    @staticmethod
    def _set_cell_shading(cell: Any, hex_color: str) -> None:
        """Set cell background color via OpenXML shading."""
        from docx.oxml import parse_xml
        from docx.oxml.ns import nsdecls

        shading_elm = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{hex_color}"/>')
        cell._tc.get_or_add_tcPr().append(shading_elm)

    @staticmethod
    def _inject_native_charts(
        output_path: str,
        charts: list[Chart],
    ) -> None:
        """Generate chart XML from PPTX and inject into the DOCX.

        Temporary PPTX files are created in the system temp directory and
        cleaned up after extraction.

        Raises:
            RenderingError: If chart generation fails for any chart.
        """
        if not charts:
            return

        chart_files: list[dict[str, object]] = []

        for idx, chart_model in enumerate(charts, start=1):
            temp_pptx: str | None = None
            try:
                # Create temp PPTX in system temp dir (not working dir)
                fd, temp_pptx = tempfile.mkstemp(
                    suffix=".pptx", prefix="reportkit_chart_"
                )
                os.close(fd)

                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                chart_data = CategoryChartData()  # type: ignore[no-untyped-call]
                chart_data.categories = chart_model.categories

                for series in chart_model.series:
                    chart_data.add_series(series.name, tuple(series.values))  # type: ignore[no-untyped-call]

                x, y, cx, cy = (
                    PPTXInches(1),
                    PPTXInches(1),
                    PPTXInches(6),
                    PPTXInches(4),
                )

                chart_type = _CHART_TYPE_MAP.get(
                    chart_model.chart_type,
                    PPTX_XL_CHART_TYPE.COLUMN_CLUSTERED,
                )

                slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
                prs.save(temp_pptx)

                with zipfile.ZipFile(temp_pptx, "r") as z:
                    excel_sheet_path: str | None = None
                    for fname in z.namelist():
                        if fname.startswith("ppt/embeddings/Microsoft_Excel_"):
                            excel_sheet_path = fname
                            break

                    if not excel_sheet_path:
                        raise RenderingError(
                            f"No embedded workbook found in chart {idx}"
                        )

                    xml = z.read("ppt/charts/chart1.xml")
                    rels = z.read("ppt/charts/_rels/chart1.xml.rels")
                    xlsx = z.read(excel_sheet_path)

                chart_files.append(
                    {
                        "idx": idx,
                        "chart_id": chart_model.chart_id,
                        "xml": xml,
                        "rels": rels,
                        "xlsx": xlsx,
                    }
                )

            except RenderingError:
                raise
            except Exception as exc:
                logger.error(
                    "Failed to generate chart %d (%s): %s",
                    idx,
                    chart_model.chart_id,
                    exc,
                )
                raise RenderingError(
                    f"Failed to generate chart '{chart_model.chart_id}'"
                ) from exc
            finally:
                if temp_pptx and os.path.exists(temp_pptx):
                    os.remove(temp_pptx)

        if chart_files:
            OXmlPatcher.inject_charts(output_path, chart_files)
